from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation(filename="TruthLens_Presentation.pptx"):
    prs = Presentation()
    
    # Slides structure
    slides_data = [
        # Slide 1
        {"title": "TruthLens", "content": "Fake News Detection Platform\n\nAnalyze • Detect • Explain\n\nA professional, portfolio-grade analytics platform powered by advanced NLP and machine learning.", "layout": 0},
        # Slide 2
        {"title": "The Problem: Misinformation", "content": "- Misinformation spreads rapidly online\n- Fake news impacts public opinion and trust\n- Manual verification is slow and scales poorly\n- Need for automated, reliable, and explainable tools", "layout": 1},
        # Slide 3
        {"title": "The Solution: TruthLens", "content": "- Automated fake news detection\n- End-to-end NLP pipeline\n- User-friendly dashboard for analytics\n- Local, privacy-first execution\n- High contrast, flat design UI for accessibility", "layout": 1},
        # Slide 4
        {"title": "Core Capabilities", "content": "1. Analyze: Ingest multi-format news feeds and curate linguistic corpora.\n2. Detect: Predict article validity using ML models.\n3. Explain: Transparent explainable AI (XAI) confidence scores.", "layout": 1},
        # Slide 5
        {"title": "Target Audience", "content": "- Journalists verifying sources\n- Researchers analyzing information spread\n- Data Scientists exploring NLP techniques\n- General users validating news articles", "layout": 1},
        # Slide 6
        {"title": "Technology Stack", "content": "- Frontend: Streamlit (Python)\n- Backend logic: Python\n- Machine Learning: Scikit-Learn\n- NLP Processing: NLTK\n- Database: SQLite\n- Visualizations: Plotly", "layout": 1},
        # Slide 7
        {"title": "Platform Workflow Pipeline", "content": "Step 1: Upload Dataset\nStep 2: Data Prep\nStep 3: Dashboard\nStep 4: Train Model\nStep 5: Detect News\nStep 6: Reports", "layout": 1},
        
        # Slide 8
        {"title": "Step 1: Upload Dataset", "content": "- Import local dataset files\n- Support for CSV, Excel (XLSX, XLS), and JSON formats\n- Up to 200MB per file\n- Maps text fields to configure the ingestion engine", "layout": 1},
        # Slide 9
        {"title": "Instant Portfolio Testing", "content": "- Don't have a dataset?\n- Populate the system with synthetic articles\n- Contains distinct vocabularies for real and fake news\n- Accelerates testing and demonstration", "layout": 1},
        
        # Slide 10
        {"title": "Step 2: Data Preparation & Cleaning", "content": "- Select ingested dataset to preprocess\n- Visual feedback on total articles vs. already preprocessed\n- Text Preprocessing and NLP Cleaning", "layout": 1},
        # Slide 11
        {"title": "NLP Cleaners Configuration", "content": "- Remove English Stopwords (e.g., 'the', 'is', 'at')\n- Apply WordNet Lemmatization (converting words to base form)\n- Apply Porter Stemming (reducing words to root form)", "layout": 1},
        # Slide 12
        {"title": "Preprocessing Dry Run", "content": "- Compares Original Text (Truncated) with Cleansed NLP Tokens\n- Ensures data quality before training\n- Transparent view into the preprocessing engine", "layout": 1},
        
        # Slide 13
        {"title": "Step 3: Analytics Dashboard", "content": "- Interactive visualizations of the curated dataset\n- Driven by Plotly\n- Fully integrated into the flat design light theme", "layout": 1},
        # Slide 14
        {"title": "Corpus Metrics", "content": "- Total Articles Ingested\n- Cleaned & Cataloged Articles\n- Unique Vocabulary Words\n- Average Words per Article", "layout": 1},
        # Slide 15
        {"title": "Class Balance & Distribution", "content": "- Donut chart showing Class Balance (Real vs. Fake)\n- Ensures balanced datasets to prevent model bias\n- Histogram showing Article Word Count Distribution", "layout": 1},
        # Slide 16
        {"title": "Word Frequency Analysis", "content": "- Filter word frequency by Real vs. Fake Articles\n- Top 15 Most Common Words visualized via Bar Chart\n- Helps identify linguistic patterns in misinformation", "layout": 1},
        
        # Slide 17
        {"title": "Step 4: Model Training", "content": "- Configure and train custom classification models\n- Select target dataset\n- Train on real-time data within the platform", "layout": 1},
        # Slide 18
        {"title": "Feature Extraction", "content": "- Vectorizer Type: TF-IDF Vectorizer\n- N-Gram Range selection (e.g., Unigrams Only)\n- Configurable Max Vocabulary Features limit", "layout": 1},
        # Slide 19
        {"title": "Classifier Algorithm Settings", "content": "- Algorithm Selection (e.g., Logistic Regression)\n- Configurable Hyperparameters (e.g., Regularization Strength C)\n- Test Set Split Ratio (%) control", "layout": 1},
        # Slide 20
        {"title": "Training Execution", "content": "- Define custom Model Version Name\n- Fit & Evaluate Model button\n- Generates models serialized safely into the local workspace", "layout": 1},
        
        # Slide 21
        {"title": "Step 5: Detect News", "content": "- The core inference engine\n- Evaluate new articles against trained models\n- Select Active Prediction Model from dropdown", "layout": 1},
        # Slide 22
        {"title": "Input Methods", "content": "- Paste Text Article directly into the text area\n- Verify Live Web URL (Future/Optional Integration)\n- Analyze & Verify Authenticity execution button", "layout": 1},
        # Slide 23
        {"title": "Detection Results", "content": "- Real-time validation outcome banners\n- Verified Real (Green) vs. Detected Fake (Red)\n- Explainable AI Confidence Scores", "layout": 1},
        # Slide 24
        {"title": "Explainable AI (XAI) Tokens", "content": "- Highlights influential tokens in the text\n- High-contrast inline spans for readability\n- Demystifies the machine learning decision process", "layout": 1},
        
        # Slide 25
        {"title": "Step 6: Reports & Logging", "content": "- System Performance Audits\n- Select Target Dataset and Target Model Run\n- Archival of historical runs", "layout": 1},
        # Slide 26
        {"title": "Unified Audit Summaries", "content": "- Maps linguistic features of the dataset\n- Records cleaning toggles used\n- Captures classification scoring details\n- Archived locally in HTML format", "layout": 1},
        
        # Slide 27
        {"title": "UI/UX Principles: Flat Design", "content": "- Platform underwent a major UI overhaul\n- Rejects 3D illusions: no drop shadows, bevels, or textures\n- Relies on hierarchy through size, color, and typography", "layout": 1},
        # Slide 28
        {"title": "Typography & Color", "content": "- Font Family: 'Outfit', sans-serif\n- High-contrast Light Mode\n- Background: Pure White | Foreground: Gray 900\n- Strict reliance on the grid", "layout": 1},
        
        # Slide 29
        {"title": "Future Enhancements", "content": "- Live URL Web Scraping capabilities\n- Advanced Deep Learning Models (Transformers/BERT)\n- Real-time API endpoints for external integrations\n- Multi-language support", "layout": 1},
        # Slide 30
        {"title": "Conclusion", "content": "Thank You!\n\nTruthLens provides an accessible, private, and powerful way to understand and combat misinformation.\n\nQuestions?", "layout": 1},
    ]
    
    for slide_info in slides_data:
        layout = prs.slide_layouts[slide_info["layout"]]
        slide = prs.slides.add_slide(layout)
        
        title_shape = slide.shapes.title
        title_shape.text = slide_info["title"]
        
        if slide_info["layout"] == 0:
            # Title slide
            subtitle = slide.placeholders[1]
            subtitle.text = slide_info["content"]
        else:
            # Bullet slide
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.text = slide_info["content"]
            
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == '__main__':
    create_presentation()
